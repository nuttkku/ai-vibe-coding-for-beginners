"""
Create PowerPoint from slide.md — Vibe Coding for Beginners
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import re
import os

# ─── Color Palette ────────────────────────────────────────────────────────────
C_NAVY      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy (title bg)
C_BLUE      = RGBColor(0x1D, 0x4E, 0xD8)   # brand blue
C_BLUE_MID  = RGBColor(0x25, 0x63, 0xEB)   # mid blue
C_CYAN      = RGBColor(0x06, 0xB6, 0xD4)   # cyan accent
C_GOLD      = RGBColor(0xF5, 0x9E, 0x0B)   # gold accent
C_GREEN     = RGBColor(0x10, 0xB9, 0x81)   # success green
C_RED       = RGBColor(0xEF, 0x44, 0x44)   # error red
C_PURPLE    = RGBColor(0x7C, 0x3A, 0xED)   # purple accent
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
C_LIGHT_BG  = RGBColor(0xF0, 0xF9, 0xFF)   # very light blue bg
C_GRAY      = RGBColor(0x64, 0x74, 0x8B)   # slate gray text
C_DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)   # dark text
C_CODE_BG   = RGBColor(0x1E, 0x1E, 0x2E)   # code block bg
C_CODE_TEXT = RGBColor(0xA6, 0xE3, 0xA1)   # code text green
C_BORDER    = RGBColor(0xBF, 0xDB, 0xFE)   # light blue border

# Slide size: 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def hex_to_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def set_font(run, name="Leelawadee UI", size=None, bold=False, color=None, italic=False):
    run.font.name = name
    if size:
        run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, font_size=18, bold=False,
                color=C_DARK_TEXT, align=PP_ALIGN.LEFT, font="Leelawadee UI",
                italic=False, wrap=True, line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, name=font, size=font_size, bold=bold, color=color, italic=italic)
    if line_spacing:
        from pptx.util import Pt
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txBox


def parse_slides(md_text):
    """Split markdown into list of (title, body_lines)."""
    raw_slides = md_text.split('\n---\n')
    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        lines = raw.split('\n')
        title = ''
        subtitle = ''
        body = []
        for line in lines:
            if line.startswith('# ') and not title:
                title = line[2:].strip()
            elif line.startswith('## ') and not subtitle:
                subtitle = line[3:].strip()
            else:
                body.append(line)
        slides.append({
            'title': title,
            'subtitle': subtitle,
            'body': [l for l in body if l.strip() or body.index(l) > 0],
        })
    return slides


def classify_slide(slide_data, idx):
    """Classify slide type."""
    title = slide_data['title']
    if idx == 0:
        return 'cover'
    if title.startswith('LAB-'):
        # Check if it's a LAB intro (has Session/ผลลัพธ์ in body)
        body_text = '\n'.join(slide_data['body'])
        if 'Session' in body_text or 'ผลลัพธ์' in body_text:
            return 'lab_intro'
        return 'lab_content'
    if 'ตาราง' in title or title.startswith('ตาราง'):
        return 'schedule'
    if 'ไอเดีย' in title:
        return 'ideas'
    if 'สรุป' in title or 'ทักษะ' in title:
        return 'summary'
    if title.startswith('"') or '70:20:10' in title:
        return 'quote'
    return 'content'


def render_cover(slide, data, prs):
    # Full navy background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_NAVY)
    # Decorative top stripe
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_CYAN)
    # Decorative bottom stripe
    add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), C_GOLD)

    # Accent block (left vertical bar)
    add_rect(slide, 0, Inches(1.8), Inches(0.12), Inches(3.5), C_CYAN)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.5), Inches(12.5), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = data['title']
    set_font(run, size=44, bold=True, color=C_WHITE)

    # Subtitle
    if data['subtitle']:
        txBox2 = slide.shapes.add_textbox(Inches(0.4), Inches(3.3), Inches(12.5), Inches(1.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = data['subtitle']
        set_font(run2, size=22, color=C_CYAN)

    # Body lines (ระยะเวลา, กลุ่มเป้าหมาย)
    y_pos = Inches(4.5)
    for line in data['body']:
        line = line.strip()
        if not line:
            continue
        txBox3 = slide.shapes.add_textbox(Inches(0.4), y_pos, Inches(12.5), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = line
        set_font(run3, size=18, color=C_GOLD)
        y_pos += Inches(0.45)

    # Logo-like badge in bottom-right
    add_rect(slide, SLIDE_W - Inches(3.2), SLIDE_H - Inches(1.6), Inches(3.0), Inches(1.3), C_BLUE)
    add_textbox(slide, SLIDE_W - Inches(3.15), SLIDE_H - Inches(1.55), Inches(2.9), Inches(1.1),
                "🤖 Vibe Coding\nWorkshop 2 วัน", font_size=14, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)


def render_lab_intro(slide, data):
    # Top colored bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.6), C_BLUE)
    # Accent line under bar
    add_rect(slide, 0, Inches(1.6), SLIDE_W, Inches(0.05), C_GOLD)

    # Title in bar
    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(12.5), Inches(1.3),
                data['title'], font_size=30, bold=True, color=C_WHITE)

    # Light background for rest
    add_rect(slide, 0, Inches(1.65), SLIDE_W, SLIDE_H - Inches(1.65), C_LIGHT_BG)

    # Render body
    y_pos = Inches(1.85)
    body_lines = [l for l in data['body'] if l.strip()]

    # Find session line, result line, and other lines
    session_line = ''
    result_line = ''
    other_lines = []
    for line in body_lines:
        if 'Session' in line and 'น.' in line:
            session_line = line.strip()
        elif line.strip().startswith('ผลลัพธ์'):
            result_line = line.strip()
        else:
            other_lines.append(line)

    if session_line:
        add_textbox(slide, Inches(0.5), y_pos, Inches(12), Inches(0.5),
                    session_line, font_size=14, color=C_GRAY, italic=True)
        y_pos += Inches(0.5)

    if result_line:
        # Result box
        add_rect(slide, Inches(0.4), y_pos, Inches(12.4), Inches(0.65), C_GREEN)
        add_textbox(slide, Inches(0.6), y_pos + Inches(0.08), Inches(12.0), Inches(0.5),
                    result_line, font_size=16, bold=True, color=C_WHITE)
        y_pos += Inches(0.8)

    for line in other_lines:
        y_pos = render_body_line(slide, line, y_pos, x_indent=Inches(0.5))
        if y_pos > Inches(7.0):
            break


def render_content(slide, data, slide_type='content'):
    """Standard content slide."""
    # Background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_WHITE)

    # Choose header color
    if slide_type == 'lab_content':
        hdr_color = C_BLUE
    elif slide_type == 'schedule':
        hdr_color = C_PURPLE
    elif slide_type == 'summary':
        hdr_color = C_GREEN
    elif slide_type == 'ideas':
        hdr_color = C_GOLD
    else:
        hdr_color = C_NAVY

    # Header bar
    hdr_h = Inches(1.25)
    add_rect(slide, 0, 0, SLIDE_W, hdr_h, hdr_color)
    # Bottom border accent
    add_rect(slide, 0, hdr_h, SLIDE_W, Inches(0.06), C_CYAN)

    # Title text
    add_textbox(slide, Inches(0.35), Inches(0.12), Inches(12.6), Inches(1.0),
                data['title'], font_size=26, bold=True, color=C_WHITE)

    # Content area background
    add_rect(slide, 0, hdr_h + Inches(0.06), SLIDE_W,
             SLIDE_H - hdr_h - Inches(0.06), C_WHITE)

    # Render body
    body_lines = data['body']
    y_pos = Inches(1.5)
    i = 0
    while i < len(body_lines):
        line = body_lines[i]
        y_pos = render_body_line(slide, line, y_pos)
        i += 1
        if y_pos > Inches(7.1):
            break


def render_quote(slide, data):
    """Full-color quote slide."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_NAVY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.1), C_GOLD)
    add_rect(slide, 0, SLIDE_H - Inches(0.1), SLIDE_W, Inches(0.1), C_GOLD)

    # Big quote mark
    add_textbox(slide, Inches(0.3), Inches(0.3), Inches(2), Inches(1.5),
                '"', font_size=100, bold=True, color=C_CYAN)

    title = data['title'].strip('"').strip('"').strip('"')
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.0),
                f'"{title}"', font_size=24, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)

    y_pos = Inches(4.8)
    for line in data['body']:
        line = line.strip()
        if not line:
            continue
        add_textbox(slide, Inches(1.0), y_pos, Inches(11.3), Inches(0.55),
                    line, font_size=16, color=C_GOLD, align=PP_ALIGN.CENTER)
        y_pos += Inches(0.55)


def render_body_line(slide, line, y_pos, x_indent=Inches(0.45)):
    """Render a single body line and return new y_pos."""
    stripped = line.strip()
    if not stripped:
        return y_pos + Inches(0.15)

    # Code-like line (docker commands, ollama commands, FROM nginx etc.)
    is_code = (stripped.startswith('docker ') or stripped.startswith('FROM ') or
               stripped.startswith('COPY ') or stripped.startswith('EXPOSE ') or
               stripped.startswith('ollama ') or stripped.startswith('wsl ') or
               stripped.startswith('http://') or
               (stripped.startswith('   ') and any(c in stripped for c in [':', '.', '/'])))

    # Section sub-header (ends with :)
    is_subheader = (stripped.endswith(':') and len(stripped) < 60 and
                    not stripped.startswith('-') and not stripped.startswith('Q:') and
                    not stripped.startswith('A:'))

    # Bullet
    if stripped.startswith('- '):
        text = stripped[2:]
        # Bullet dot
        add_rect(slide, x_indent, y_pos + Inches(0.15), Inches(0.12), Inches(0.12), C_CYAN)
        txBox = slide.shapes.add_textbox(x_indent + Inches(0.22), y_pos,
                                          SLIDE_W - x_indent - Inches(0.6), Inches(0.48))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        set_font(run, size=17, color=C_DARK_TEXT)
        return y_pos + Inches(0.42)

    elif is_code:
        # Code box
        line_h = Inches(0.42)
        add_rect(slide, x_indent, y_pos, SLIDE_W - x_indent - Inches(0.3),
                 line_h, C_CODE_BG)
        add_textbox(slide, x_indent + Inches(0.15), y_pos + Inches(0.03),
                    SLIDE_W - x_indent - Inches(0.7), line_h,
                    stripped, font_size=14, color=C_CODE_TEXT,
                    font="Courier New")
        return y_pos + line_h + Inches(0.08)

    elif is_subheader:
        # Bold sub-header with colored underline
        add_textbox(slide, x_indent, y_pos, SLIDE_W - x_indent - Inches(0.3), Inches(0.48),
                    stripped, font_size=18, bold=True, color=C_BLUE_MID)
        add_rect(slide, x_indent, y_pos + Inches(0.42), Inches(1.5), Inches(0.03), C_CYAN)
        return y_pos + Inches(0.52)

    elif stripped.startswith('Q:'):
        # FAQ question
        add_rect(slide, x_indent, y_pos, SLIDE_W - x_indent - Inches(0.3), Inches(0.48), C_BLUE)
        add_textbox(slide, x_indent + Inches(0.1), y_pos + Inches(0.04),
                    SLIDE_W - x_indent - Inches(0.5), Inches(0.4),
                    stripped, font_size=15, bold=True, color=C_WHITE)
        return y_pos + Inches(0.55)

    elif stripped.startswith('A:'):
        # FAQ answer
        add_textbox(slide, x_indent + Inches(0.2), y_pos, SLIDE_W - x_indent - Inches(0.5), Inches(0.48),
                    stripped, font_size=15, color=C_DARK_TEXT, italic=True)
        return y_pos + Inches(0.5)

    elif ' — ' in stripped and not stripped.startswith('#'):
        # Key — Value pair
        parts = stripped.split(' — ', 1)
        key_txt = parts[0].strip()
        val_txt = parts[1].strip() if len(parts) > 1 else ''
        add_textbox(slide, x_indent, y_pos, Inches(3.5), Inches(0.48),
                    key_txt, font_size=16, bold=True, color=C_BLUE_MID)
        add_textbox(slide, x_indent + Inches(3.6), y_pos, Inches(9.0), Inches(0.48),
                    val_txt, font_size=16, color=C_DARK_TEXT)
        return y_pos + Inches(0.44)

    else:
        # Normal paragraph text
        h = Inches(0.48) if len(stripped) < 80 else Inches(0.72)
        add_textbox(slide, x_indent, y_pos, SLIDE_W - x_indent - Inches(0.3), h,
                    stripped, font_size=17, color=C_DARK_TEXT)
        return y_pos + h + Inches(0.04)


def build_pptx(md_path, out_path):
    with open(md_path, encoding='utf-8') as f:
        md_text = f.read()

    slides_data = parse_slides(md_text)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # completely blank

    for idx, data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        stype = classify_slide(data, idx)

        if stype == 'cover':
            render_cover(slide, data, prs)
        elif stype == 'lab_intro':
            render_lab_intro(slide, data)
        elif stype == 'quote':
            render_quote(slide, data)
        else:
            render_content(slide, data, stype)

        # Slide number (except cover)
        if idx > 0:
            add_textbox(slide, SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.4),
                        Inches(0.8), Inches(0.35),
                        str(idx + 1), font_size=11, color=C_GRAY,
                        align=PP_ALIGN.RIGHT)

    prs.save(out_path)
    print(f"Saved: {out_path}  ({len(slides_data)} slides)")


if __name__ == '__main__':
    base = os.path.dirname(os.path.abspath(__file__))
    md_path = os.path.join(base, 'slide.md')
    out_path = os.path.join(base, 'vibe-coding-slides.pptx')
    build_pptx(md_path, out_path)
